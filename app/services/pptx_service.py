"""
PPTX Service — US2, US3
Gera arquivos .pptx a partir de slides gerados pela IA.
Suporta templates de marca (.pptx base) para manter identidade visual (US3).
"""
import uuid
import shutil
import re
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.config import settings
from app.schemas.presentation import SlideContent

# ── Caminhos ──────────────────────────────────────────────────────────────────

# Pasta de templates: raiz_do_projeto/templates/
# No Docker: /app/templates/ (mapeado no docker-compose como volume)
# Localmente: ./templates/ (relativo à raiz do projeto)
TEMPLATES_DIR = Path(__file__).resolve().parent.parent.parent / "templates"

# Pasta de saída para os .pptx gerados
STORAGE_PATH = Path(settings.LOCAL_STORAGE_PATH)


# ── Helpers internos ──────────────────────────────────────────────────────────

def _ensure_storage() -> None:
    """Garante que a pasta de outputs existe."""
    STORAGE_PATH.mkdir(parents=True, exist_ok=True)


def _load_template(template_name: Optional[str]) -> Optional[Path]:
    """
    US3 — Tenta localizar o arquivo .pptx de template pelo nome.
    Retorna o Path se encontrado, None caso contrário (usa tema padrão).

    Como usar:
      1. Crie a pasta  templates/  na raiz do projeto (mesmo nível de app/)
      2. Coloque o arquivo base da empresa:  templates/minha_empresa.pptx
      3. Passe template_name="minha_empresa" na requisição

    O arquivo .pptx base deve ter pelo menos:
      - Slide layout 0: capa (título + subtítulo)
      - Slide layout 1: conteúdo (título + área de texto com bullets)
    """
    if not template_name or template_name.strip().lower() == "default":
        return None

    candidates = [
        TEMPLATES_DIR / f"{template_name}.pptx",
        TEMPLATES_DIR / template_name,               # já com extensão
        TEMPLATES_DIR / f"{template_name}.pptx".lower(),
    ]
    for path in candidates:
        if path.exists() and path.suffix.lower() == ".pptx":
            return path

    return None


def _available_templates() -> List[str]:
    """Retorna os nomes dos templates disponíveis (sem extensão)."""
    if not TEMPLATES_DIR.exists():
        return []
    return [p.stem for p in TEMPLATES_DIR.glob("*.pptx")]


def _safe_placeholder(slide, index: int):
    """Retorna o placeholder se existir, None caso contrário."""
    try:
        return slide.placeholders[index]
    except (KeyError, IndexError):
        return None


# ── Construtores de slides ─────────────────────────────────────────────────────

def _add_title_slide(prs: Presentation, title: str) -> None:
    """Slide de capa usando layout 0 do template (ou padrão)."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    title_ph = _safe_placeholder(slide, 0)
    if title_ph:
        title_ph.text = title

    subtitle_ph = _safe_placeholder(slide, 1)
    if subtitle_ph:
        subtitle_ph.text = "Gerado automaticamente — Plataforma de Documentação Inteligente"


def _add_content_slide(prs: Presentation, slide_data: SlideContent) -> None:
    """
    Slide de conteúdo usando layout 1 do template.
    Adiciona título, bullets, notas do apresentador e bloco de código (se houver).
    """
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # Título
    title_ph = _safe_placeholder(slide, 0)
    if title_ph:
        title_ph.text = slide_data.title

    # Bullets
    body_ph = _safe_placeholder(slide, 1)
    if body_ph and slide_data.bullets:
        tf = body_ph.text_frame
        tf.clear()
        for i, bullet in enumerate(slide_data.bullets):
            para = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            para.text = bullet
            para.level = 0

    # Bloco de código — caixa de texto flutuante estilo dark (US4/US9)
    if slide_data.code_snippet:
        _add_code_box(slide, slide_data.code_snippet)

    # Notas do apresentador
    if slide_data.notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = slide_data.notes


def _add_code_box(slide, code: str) -> None:
    """Adiciona uma caixa de texto com estilo de editor de código e syntax highlighting."""
    left, top, width, height = Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.6)
    box = slide.shapes.add_textbox(left, top, width, height)

    fill = box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)

    tf = box.text_frame
    tf.word_wrap = True

    pattern = re.compile(
        r'(?P<string>"(?:\\.|[^"\\])*"|\'(?:\\.|[^\'\\])*\')|'
        r'(?P<comment>#.*?$|//.*?$)|'
        r'(?P<keyword>\b(?:def|class|import|from|return|if|else|elif|for|while|try|except|async|await|with|as|pass|yield|break|continue|True|False|None|and|or|not|in|is|const|let|var|function|export|default|type|interface)\b)|'
        r'(?P<other>(?:(?!["\']|#|//|\b(?:def|class|import|from|return|if|else|elif|for|while|try|except|async|await|with|as|pass|yield|break|continue|True|False|None|and|or|not|in|is|const|let|var|function|export|default|type|interface)\b).)+)'
    )

    # Separa por linhas para forçar a quebra correta no PPTX
    for line_idx, line in enumerate(code.split("\n")):
        para = tf.paragraphs[0] if line_idx == 0 else tf.add_paragraph()
        # Remove espaçamento inútil entre as linhas de código
        para.space_before = Pt(0)
        para.space_after = Pt(0)

        if not line.strip():
            run = para.add_run()
            run.text = " "
            continue

        for match in pattern.finditer(line):
            run = para.add_run()
            run.text = match.group(0)
            run.font.size, run.font.name, run.font.bold = Pt(9), "Cascadia Code", False

            if match.group("string"):
                run.font.color.rgb = RGBColor(0xA6, 0xE3, 0xA1)
            elif match.group("comment"):
                run.font.color.rgb = RGBColor(0x7F, 0x84, 0x8E)
            elif match.group("keyword"):
                run.font.color.rgb = RGBColor(0xC6, 0x78, 0xDD)
            else:
                run.font.color.rgb = RGBColor(0xDC, 0xDC, 0xDC)

# ── API pública ───────────────────────────────────────────────────────────────

def export_to_pptx(
    slides: List[SlideContent],
    title: str,
    template_name: Optional[str] = None,
) -> tuple[str, str]:
    """
    US2, US3 — Gera e salva o arquivo .pptx.

    Args:
        slides:        Lista de slides gerados pela IA.
        title:         Título da apresentação (usado no slide de capa).
        template_name: Nome do arquivo .pptx em templates/ (sem extensão).
                       Passar None ou "default" usa o tema branco padrão do python-pptx.

    Returns:
        (filepath_absoluto, filename) — o filename é usado para montar a URL de download.
    """
    _ensure_storage()

    template_path = _load_template(template_name)

    if template_path:
        # US3: copia o template para não modificar o original
        tmp_copy = STORAGE_PATH / f"_tpl_{uuid.uuid4().hex}.pptx"
        shutil.copy2(template_path, tmp_copy)
        prs = Presentation(str(tmp_copy))
        tmp_copy.unlink(missing_ok=True)
    else:
        prs = Presentation()
        # Proporção widescreen 16:9
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

    _add_title_slide(prs, title)

    for slide_data in slides:
        _add_content_slide(prs, slide_data)

    filename = f"{uuid.uuid4().hex}.pptx"
    filepath = STORAGE_PATH / filename
    prs.save(str(filepath))

    return str(filepath), filename


def list_templates() -> List[str]:
    """Retorna os templates disponíveis na pasta templates/."""
    return _available_templates()
